# -*- coding: utf-8 -*-
"""构建课件 PPTX 的小工具库。

设计目标：
  * 16:9 版面，中文字体（微软雅黑）在 Windows / Office 下开箱即用；
  * 用一组简单的"幻灯片描述"元组来写内容，与排版代码解耦；
  * 代码页按行数自动缩放字号，尽量避免溢出。

幻灯片描述（见各 content/wNN.py）支持的类型：
  ("title",   主标题, 副标题)
  ("section", 编号, 章节标题, 可选小字)
  ("bullets", 标题, [条目, ...])            条目以 "- " 前缀表示次级
  ("code",    标题, 代码字符串, 说明)
  ("table",   标题, [[表头...], [行...], ...], 说明)
  ("two",     标题, 左标题, [左条目...], 右标题, [右条目...])
  ("key",     标题, 要点正文)               整页强调一句话
  ("ascii",   标题, 等宽示意图, 说明)
"""

import math
import re

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- 主题配置
EA_FONT = '微软雅黑'          # 中文字体
LATIN_FONT = 'Segoe UI'       # 西文字体
MONO_FONT = 'Consolas'        # 代码字体

NAVY = RGBColor(0x12, 0x39, 0x5B)      # 主色：深蓝
ACCENT = RGBColor(0xE0, 0x7B, 0x39)    # 强调色：橙
INK = RGBColor(0x24, 0x2A, 0x33)       # 正文
MUTED = RGBColor(0x6B, 0x75, 0x82)     # 次要文字
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)     # 浅底
CODE_BG = RGBColor(0xF6, 0xF7, 0xF9)
RULE = RGBColor(0xD6, 0xDC, 0xE4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.62)
BODY_TOP = Inches(1.42)
BODY_H = Inches(5.32)
BODY_W = SLIDE_W - 2 * MARGIN


# ---------------------------------------------------------------- 字体工具
def _style_run(run, size, bold=False, color=INK, mono=False, italic=False):
    """设置一个 run 的字体。python-pptx 只写 a:latin，中文需手工补 a:ea。"""
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = MONO_FONT if mono else LATIN_FONT
    rPr = f._rPr
    for tag in ('a:ea', 'a:cs'):
        for old in rPr.findall(qn(tag)):
            rPr.remove(old)
    latin = rPr.find(qn('a:latin'))
    ea = rPr.makeelement(qn('a:ea'), {'typeface': MONO_FONT if mono else EA_FONT})
    if latin is not None:
        latin.addnext(ea)
    else:
        rPr.append(ea)


def _textbox(slide, left, top, width, height, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    # 固定尺寸：否则不换行的文本框会被渲染器自动居中，左边距忽宽忽窄
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    return box, tf


_SEGMENT = re.compile(r'(\*\*[^*]+\*\*|`[^`]+`)')


def _plain(text):
    """去掉 **强调** 与 `等宽` 标记，用于度量文字宽度。"""
    return text.replace('**', '').replace('`', '')


def _add_runs(p, text, size, bold, color, mono):
    """把一行文本按 **强调** / `等宽` 切成多个 run。

    代码/示意图（mono）原样输出：Python 的 ** 幂运算符不能被当成强调标记。
    """
    if mono:
        run = p.add_run()
        run.text = text
        _style_run(run, size, bold, color, True)
        return
    for seg in _SEGMENT.split(text):
        if not seg:
            continue
        if seg.startswith('**') and seg.endswith('**'):
            run = p.add_run()
            run.text = seg[2:-2]
            _style_run(run, size, True, NAVY if color is INK else color, mono)
        elif seg.startswith('`') and seg.endswith('`'):
            run = p.add_run()
            run.text = seg[1:-1]
            _style_run(run, size * 0.94, bold, color, True)
        else:
            run = p.add_run()
            run.text = seg
            _style_run(run, size, bold, color, mono)


def _para(tf, text, size, bold=False, color=INK, mono=False,
          space_before=0, space_after=6, first=False, align=PP_ALIGN.LEFT,
          line_spacing=None, indent_level=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if line_spacing:
        p.line_spacing = line_spacing
    p.level = indent_level
    _add_runs(p, text, size, bold, color, mono)
    return p


def _rect(slide, left, top, width, height, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


# ---------------------------------------------------------------- 排版度量
def _ems(text):
    """估算一段文字的宽度，单位为 em（1 em = 当前字号）。

    CJK 与全角标点按 1.0 计，ASCII 按 0.55 计 —— 用于估算换行与列宽。
    """
    w = 0.0
    for ch in _plain(text):
        w += 1.0 if ord(ch) > 0x2E80 else 0.55
    return w


def _wrapped_lines(text, size_pt, width_pt):
    """给定字号与可用宽度，估算这段文字会占几行。"""
    if not text:
        return 1
    return max(1, math.ceil(_ems(text) * size_pt / width_pt))


def _fit_size(items, width_pt, height_pt, hi, lo, line_spacing, gap_ratio):
    """在 [lo, hi] 中挑最大的字号，使全部条目仍能装进 height_pt。

    items 为 (文本, 相对字号比例) 的列表。
    """
    size = hi
    while size > lo:
        total = 0.0
        for text, ratio in items:
            s = size * ratio
            total += _wrapped_lines(text, s, width_pt) * s * line_spacing
            total += s * gap_ratio
        if total <= height_pt:
            break
        size -= 0.5
    return size


# ---------------------------------------------------------------- 版面构件
def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _slide_header(slide, title):
    """页眉：标题 + 橙色短线。"""
    _, tf = _textbox(slide, MARGIN, Inches(0.46), BODY_W, Inches(0.6))
    _para(tf, title, 26, bold=True, color=NAVY, first=True, space_after=0)
    _rect(slide, MARGIN, Inches(1.16), Inches(1.05), Pt(3.2), fill=ACCENT)


def _footer(slide, label, number):
    _, tf = _textbox(slide, MARGIN, Inches(6.92), Inches(9.0), Inches(0.34))
    _para(tf, label, 10, color=MUTED, first=True, space_after=0)
    _, tf2 = _textbox(slide, SLIDE_W - MARGIN - Inches(1.0), Inches(6.92),
                      Inches(1.0), Inches(0.34))
    _para(tf2, str(number), 10, color=MUTED, first=True, space_after=0,
          align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------- 各类页面
def _add_title(prs, main, sub, meta):
    slide = _blank(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
    _rect(slide, 0, Inches(4.05), SLIDE_W, Pt(3.2), fill=ACCENT)
    _, tf = _textbox(slide, MARGIN, Inches(2.30), SLIDE_W - 2 * MARGIN, Inches(1.7))
    _para(tf, main, 40, bold=True, color=WHITE, first=True, space_after=10)
    if sub:
        _para(tf, sub, 19, color=RGBColor(0xC7, 0xD3, 0xE2), space_after=0)
    _, tf2 = _textbox(slide, MARGIN, Inches(4.55), SLIDE_W - 2 * MARGIN, Inches(1.6))
    firstline = True
    for line in meta:
        _para(tf2, line, 13, color=RGBColor(0xA8, 0xB8, 0xCC),
              first=firstline, space_after=5)
        firstline = False
    return slide


def _add_section(prs, number, title, note=''):
    slide = _blank(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT)
    _rect(slide, 0, Inches(2.55), Inches(0.30), Inches(1.9), fill=ACCENT)
    _, tf = _textbox(slide, Inches(0.95), Inches(2.62), Inches(11.4), Inches(1.9))
    _para(tf, number, 15, bold=True, color=ACCENT, first=True, space_after=8)
    _para(tf, title, 34, bold=True, color=NAVY, space_after=8)
    if note:
        _para(tf, note, 14, color=MUTED, space_after=0)
    return slide


def _add_bullets(prs, title, items):
    slide = _blank(prs)
    _slide_header(slide, title)
    body_w_pt = BODY_W / 12700.0
    body_h_pt = BODY_H / 12700.0

    parsed = []
    for raw in items:
        sub = raw.startswith('- ')
        text = raw[2:] if sub else raw
        head = (text.startswith('**') and text.endswith('**')
                and text.count('**') == 2)
        if head:
            text = text[2:-2]
        parsed.append((text, sub, head))

    size = _fit_size([(t, 0.82 if sub else 1.0) for t, sub, _ in parsed],
                     body_w_pt, body_h_pt * 0.92, hi=20, lo=11.5,
                     line_spacing=1.32, gap_ratio=0.62)
    gap = size * 0.62

    _, tf = _textbox(slide, MARGIN, BODY_TOP, BODY_W, BODY_H)
    first = True
    for text, sub, bold in parsed:
        if sub:
            _para(tf, '·  ' + text, size * 0.82, color=MUTED, first=first,
                  space_after=gap * 0.72, indent_level=1, line_spacing=1.3)
        else:
            _para(tf, ('▍ ' if bold else '•  ') + text, size, bold=bold,
                  color=NAVY if bold else INK, first=first,
                  space_after=gap, line_spacing=1.32)
        first = False
    return slide


def _add_code(prs, title, code, caption=''):
    slide = _blank(prs)
    _slide_header(slide, title)
    lines = code.rstrip('\n').split('\n')
    n = len(lines)
    widest = max((_ems(l) for l in lines), default=1)

    cap_h = Inches(0.46) if caption else Inches(0)
    avail_h_pt = (BODY_H - cap_h) / 12700.0
    avail_w_pt = (BODY_W - Inches(0.72)) / 12700.0

    # 行距用绝对磅值，保证盒子高度可精确计算，不会溢出
    LEAD = 1.30                                   # 行高 / 字号
    PAD = 22.0                                    # 盒子上下内边距合计（pt）
    by_rows = (avail_h_pt - PAD) / (n * LEAD)
    by_cols = avail_w_pt / max(widest, 1) / 0.60  # 等宽字体约 0.60 em/字符
    size = max(8.5, min(15.5, by_rows, by_cols))

    box_h_pt = min(float(BODY_H) / 12700.0 - float(cap_h) / 12700.0,
                   n * size * LEAD + PAD)
    box_h = Emu(int(box_h_pt * 12700))
    _rect(slide, MARGIN, BODY_TOP, BODY_W, box_h, fill=CODE_BG, line=RULE)
    _, tf = _textbox(slide, MARGIN + Inches(0.28), BODY_TOP + Inches(0.11),
                     BODY_W - Inches(0.56), box_h - Inches(0.22))
    first = True
    for line in lines:
        _para(tf, line if line else ' ', size, mono=True, color=INK,
              first=first, space_after=0, line_spacing=Pt(size * LEAD))
        first = False
    if caption:
        _, tf2 = _textbox(slide, MARGIN, BODY_TOP + box_h + Inches(0.13),
                          BODY_W, Inches(0.36))
        _para(tf2, caption, 13, color=MUTED, first=True, space_after=0)
    return slide


def _add_ascii(prs, title, art, caption=''):
    slide = _blank(prs)
    _slide_header(slide, title)
    lines = art.rstrip('\n').split('\n')
    n = len(lines)
    widest = max((_ems(l) for l in lines), default=1)
    cap_h = Inches(0.46) if caption else Inches(0)
    avail_h_pt = (BODY_H - cap_h) / 12700.0
    avail_w_pt = BODY_W / 12700.0

    LEAD = 1.30
    by_rows = avail_h_pt / (n * LEAD)
    by_cols = avail_w_pt / max(widest, 1) / 0.60
    size = max(8.5, min(17.0, by_rows, by_cols))

    art_h_pt = n * size * LEAD
    top = BODY_TOP + Emu(int(max(0.0, (avail_h_pt - art_h_pt) / 2) * 12700))
    _, tf = _textbox(slide, MARGIN, top, BODY_W,
                     Emu(int(art_h_pt * 12700)))
    first = True
    for line in lines:
        _para(tf, line if line else ' ', size, mono=True, color=NAVY,
              first=first, space_after=0, line_spacing=Pt(size * LEAD))
        first = False
    if caption:
        _, tf2 = _textbox(slide, MARGIN, Inches(6.44), BODY_W, Inches(0.36))
        _para(tf2, caption, 13, color=MUTED, first=True, space_after=0)
    return slide


def _add_table(prs, title, rows, caption=''):
    slide = _blank(prs)
    _slide_header(slide, title)
    nrow, ncol = len(rows), len(rows[0])
    cap_h = Inches(0.62) if caption else Inches(0)
    avail_h_pt = (BODY_H - cap_h) / 12700.0 * 0.97    # 留少量余量
    total_w_pt = BODY_W / 12700.0

    # 列宽按各列最长内容分配，并做上下限收敛，避免某一列被挤到换行
    weights = []
    for c in range(ncol):
        w = max(_ems(str(rows[r][c])) for r in range(nrow))
        weights.append(min(max(w, 4.0), 34.0))
    tot = sum(weights)
    col_w = [total_w_pt * w / tot for w in weights]

    # 在能装下的前提下取最大字号
    size = 17.0
    while size > 9.0:
        pad = 16.0                      # 单元格左右内边距（pt）
        h = 0.0
        for r in range(nrow):
            ls = max(_wrapped_lines(str(rows[r][c]), size, col_w[c] - pad)
                     for c in range(ncol))
            h += max(ls * size * 1.30 + 12, 34.6)
        if h <= avail_h_pt:
            break
        size -= 0.5

    row_h = []
    for r in range(nrow):
        ls = max(_wrapped_lines(str(rows[r][c]), size, col_w[c] - 16.0)
                 for c in range(ncol))
        row_h.append(max(ls * size * 1.30 + 12, 34.6))
    height = Emu(int(sum(row_h) * 12700))

    shape = slide.shapes.add_table(nrow, ncol, MARGIN, BODY_TOP, BODY_W, height)
    table = shape.table
    for c in range(ncol):
        table.columns[c].width = Emu(int(col_w[c] * 12700))
    for r, row in enumerate(rows):
        table.rows[r].height = Emu(int(row_h[r] * 12700))
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.margin_left = cell.margin_right = Inches(0.11)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if r == 0 else (
                WHITE if r % 2 else LIGHT)
            tf = cell.text_frame
            tf.word_wrap = True
            _para(tf, str(val), size, bold=(r == 0),
                  color=WHITE if r == 0 else INK, first=True, space_after=0,
                  line_spacing=1.22)
    if caption:
        # 固定在版心底部，而不是紧跟表格 —— 表格实际行高可能略高于估算
        top = min(BODY_TOP + height + Inches(0.16), Inches(6.48))
        _, tf2 = _textbox(slide, MARGIN, top, BODY_W, Inches(0.36))
        _para(tf2, caption, 13, color=MUTED, first=True, space_after=0)
    return slide


def _add_two(prs, title, lhead, litems, rhead, ritems):
    slide = _blank(prs)
    _slide_header(slide, title)
    colw = (BODY_W - Inches(0.5)) / 2
    for i, (head, items) in enumerate(((lhead, litems), (rhead, ritems))):
        left = MARGIN + Emu(int(i * (colw + Inches(0.5))))
        _rect(slide, left, BODY_TOP, Emu(int(colw)), Inches(0.46), fill=NAVY)
        _, th = _textbox(slide, left + Inches(0.16), BODY_TOP + Inches(0.09),
                         Emu(int(colw)) - Inches(0.32), Inches(0.34))
        _para(th, head, 16, bold=True, color=WHITE, first=True, space_after=0)
        n = len(items)
        size = 16 if n <= 5 else (14.5 if n <= 7 else 13)
        _, tf = _textbox(slide, left + Inches(0.10), BODY_TOP + Inches(0.66),
                         Emu(int(colw)) - Inches(0.20), BODY_H - Inches(0.8))
        first = True
        for raw in items:
            mono = raw.startswith('`') and raw.endswith('`')
            text = raw.strip('`')
            _para(tf, ('· ' if not mono else '') + text, size, mono=mono,
                  color=INK, first=first, space_after=8, line_spacing=1.25)
            first = False
    return slide


def _add_key(prs, title, text):
    slide = _blank(prs)
    _slide_header(slide, title)
    _rect(slide, MARGIN, Inches(2.28), BODY_W, Inches(2.5), fill=LIGHT)
    _rect(slide, MARGIN, Inches(2.28), Inches(0.09), Inches(2.5), fill=ACCENT)
    _, tf = _textbox(slide, MARGIN + Inches(0.55), Inches(2.60),
                     BODY_W - Inches(1.1), Inches(1.9))
    size = 26 if len(text) <= 40 else (22 if len(text) <= 70 else 18)
    _para(tf, text, size, bold=True, color=NAVY, first=True, space_after=0,
          line_spacing=1.45)
    return slide


# ---------------------------------------------------------------- 构建入口
_BUILDERS = {
    'section': lambda prs, s: _add_section(prs, s[1], s[2],
                                           s[3] if len(s) > 3 else ''),
    'bullets': lambda prs, s: _add_bullets(prs, s[1], s[2]),
    'code': lambda prs, s: _add_code(prs, s[1], s[2],
                                     s[3] if len(s) > 3 else ''),
    'ascii': lambda prs, s: _add_ascii(prs, s[1], s[2],
                                       s[3] if len(s) > 3 else ''),
    'table': lambda prs, s: _add_table(prs, s[1], s[2],
                                       s[3] if len(s) > 3 else ''),
    'two': lambda prs, s: _add_two(prs, s[1], s[2], s[3], s[4], s[5]),
    'key': lambda prs, s: _add_key(prs, s[1], s[2]),
}


def build(meta, slides, out_path):
    """meta = {'title','subtitle','footer','info':[...]}；slides 见模块文档。"""
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    _add_title(prs, meta['title'], meta.get('subtitle', ''), meta.get('info', []))
    for i, spec in enumerate(slides, start=2):
        kind = spec[0]
        if kind == 'title':
            continue
        slide = _BUILDERS[kind](prs, spec)
        if kind != 'section':
            _footer(slide, meta.get('footer', ''), i)
    prs.save(out_path)
    return len(prs.slides)
